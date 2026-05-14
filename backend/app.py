# app.py
# Plagiarism backend — updated with external site checks, Gemini support, pie charts, and "check all" endpoint.

from flask import Flask, request, jsonify, make_response
from flask_cors import CORS
from pymongo import MongoClient
from bson.objectid import ObjectId
from datetime import datetime
import requests
import os
from io import BytesIO
import re
import json
import time
import logging
import base64
from html import unescape

# Optional plotting imports (matplotlib). We'll import inside a try/except.
try:
    import matplotlib.pyplot as plt
    MATPLOTLIB_AVAILABLE = True
except Exception:
    MATPLOTLIB_AVAILABLE = False

time.sleep(0.05)

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("plagiarism-backend")

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

# ---------- MongoDB config ----------
MONGO_URI = os.environ.get(
    "MONGO_URI",
    "mongodb+srv://Jahnavi:jahnavi@cluster0.3rq0cv5.mongodb.net/plagiarism_db?retryWrites=true&w=majority"
)
client = MongoClient(MONGO_URI)
db = client.plagiarism_db

users_collection = db.users
assignments_collection = db.assignments
submissions_collection = db.submissions

try:
    users_collection.create_index("email", unique=True)
except Exception:
    pass

# ---------- Helpers ----------
def serialize_doc(doc):
    if doc:
        doc['_id'] = str(doc['_id'])
        if 'createdAt' in doc and isinstance(doc['createdAt'], datetime):
            doc['createdAt'] = doc['createdAt'].isoformat()
        if 'deadline' in doc and isinstance(doc['deadline'], datetime):
            doc['deadline'] = doc['deadline'].isoformat()
        if 'submissionDate' in doc and isinstance(doc['submissionDate'], datetime):
            doc['submissionDate'] = doc['submissionDate'].isoformat()
    return doc

def get_user_by_id(user_id):
    try:
        user = users_collection.find_one({"_id": ObjectId(user_id)})
        if user:
            return serialize_doc(user)
    except Exception as e:
        logger.exception(f"Error fetching user by id {user_id}: {e}")
    return None

def extract_text_from_docx(file_stream):
    try:
        from docx import Document
        doc = Document(file_stream)
        return '\n'.join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"Error processing DOCX: {e}"

def extract_text_from_pdf(file_stream):
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in reader.pages:
            text += (page.extract_text() or "") + '\n\n'
        return text
    except Exception as e:
        return f"Error processing PDF: {e}"

def extract_text_from_pptx(file_stream):
    try:
        import pptx
        prs = pptx.Presentation(file_stream)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    parts.append(para.text)
        return "\n".join(parts)
    except Exception as e:
        return f"Error processing PPTX: {e}"

def safe_sentence_split(text):
    if not text: return []
    parts = re.split(r'(?<=[\.\?\!])\s+', text.strip())
    return [p.strip() for p in parts if p.strip()]

def word_set(s):
    return set(re.findall(r'\w+', (s or '').lower()))

def compute_similarity_percent(a_text, b_text):
    a_words = word_set(a_text)
    b_words = word_set(b_text)
    if not a_words: return 0
    inter = a_words.intersection(b_words)
    perc = round((len(inter) / max(1, len(a_words))) * 100)
    return min(100, perc)

def generate_pie_image(pie_data, pie_colors):
    labels = list(pie_data.keys())
    sizes = [max(0, float(pie_data.get(k, 0))) for k in labels]
    colors = [pie_colors.get(k, '#cccccc') for k in labels]

    if MATPLOTLIB_AVAILABLE:
        try:
            fig, ax = plt.subplots(figsize=(6,4), dpi=100)
            def autopct(pct):
                return ('%1.0f%%' % pct) if pct >= 1 else ''
            ax.pie(sizes if any(sizes) else [1], labels=labels if any(sizes) else ['No Data'],
                   autopct=autopct, startangle=90, colors=colors, wedgeprops=dict(width=0.5))
            ax.axis('equal')
            buf = BytesIO()
            plt.tight_layout()
            fig.savefig(buf, format='png', bbox_inches='tight', transparent=False)
            plt.close(fig)
            buf.seek(0)
            img_bytes = buf.getvalue()
            encoded = base64.b64encode(img_bytes).decode('utf-8')
            return f"data:image/png;base64,{encoded}"
        except Exception:
            logger.exception("matplotlib pie generation failed, falling back to SVG")

    try:
        total = sum(sizes) or 1
        cx, cy, r = 120, 120, 100
        svg_parts = [f'<svg xmlns="http://www.w3.org/2000/svg" width="300" height="260" viewBox="0 0 300 260">']
        angle = 0.0
        import math
        for i, (lbl, val) in enumerate(zip(labels, sizes)):
            frac = val / total if total else 0
            a1 = angle
            a2 = angle + frac * 360.0
            x1 = cx + r * math.cos(math.radians(a1 - 90))
            y1 = cy + r * math.sin(math.radians(a1 - 90))
            x2 = cx + r * math.cos(math.radians(a2 - 90))
            y2 = cy + r * math.sin(math.radians(a2 - 90))
            large = 1 if (a2 - a1) > 180 else 0
            color = colors[i]
            path = f'M {cx} {cy} L {x1} {y1} A {r} {r} 0 {large} 1 {x2} {y2} Z'
            svg_parts.append(f'<path d="{path}" fill="{color}" stroke="#ffffff" stroke-width="1"/>')
            angle = a2
        lx = 10; ly = 10
        for i, lbl in enumerate(labels):
            svg_parts.append(f'<rect x="{lx}" y="{ly + i*20}" width="12" height="12" fill="{colors[i]}"/>')
            svg_parts.append(f'<text x="{lx+18}" y="{ly + i*20 + 11}" font-size="12" fill="#333">{lbl} ({sizes[i]})</text>')
        svg_parts.append('</svg>')
        svg = ''.join(svg_parts)
        encoded = base64.b64encode(svg.encode('utf-8')).decode('utf-8')
        return f"data:image/svg+xml;base64,{encoded}"
    except Exception:
        logger.exception("SVG pie generation failed")
        return None

def summarize_matches_to_ids(local_matches, external_matches):
    local_summary = {}
    for m in local_matches:
        src = m.get('sourceId', '')
        # normalize roll id extraction
        roll = str(src).replace('Roll No.', '').replace('RollNo', '').replace('Roll', '').strip()
        try:
            pct = int(round(float(m.get('sim', 0))))
        except Exception:
            pct = 0
        if roll:
            local_summary[roll] = max(local_summary.get(roll, 0), pct)
    external_summary = {}
    for m in external_matches:
        src = m.get('sourceId', '').strip()
        try:
            pct = int(round(float(m.get('sim', 0))))
        except Exception:
            pct = 0
        if src:
            existing = external_summary.get(src, 0)
            external_summary[src] = max(existing, pct)
    return local_summary, external_summary

# Local heuristic analyzer (returns phrases for highlighting)
def analyze_with_local_heuristic(submitted_content, other_submissions_list, external_sources):
    structured = {'localMatches': [], 'externalMatches': [], 'aiMatches': []}
    sentences = safe_sentence_split(submitted_content)
    local_total = 0; external_total = 0; ai_total = 0

    for sent in sentences:
        if len(sent) < 10: continue
        best_local = None
        for s in other_submissions_list:
            sim = compute_similarity_percent(sent, s.get('content',''))
            if sim >= 15:
                if not best_local or sim > best_local['sim']:
                    best_local = {'phrase': sent, 'sim': sim, 'sourceId': f"Roll No. {s.get('roll','N/A')}"}
        if best_local:
            structured['localMatches'].append(best_local)
            local_total += best_local['sim']
            continue

        if isinstance(external_sources, dict):
            for name, ext_text in external_sources.items():
                sim = compute_similarity_percent(sent, ext_text)
                if sim >= 12:
                    structured['externalMatches'].append({'phrase': sent, 'sim': sim, 'sourceId': name})
                    external_total += sim
                    break
        else:
            sim = compute_similarity_percent(sent, external_sources or "")
            if sim >= 12:
                structured['externalMatches'].append({'phrase': sent, 'sim': sim, 'sourceId': 'External Source'})
                external_total += sim

    total_sentences = max(1, len(sentences))
    overall_score = min(100, round((local_total + external_total + ai_total) / total_sentences))
    structured['overallSimilarity'] = overall_score

    pie_data = {
        "Local (Students)": sum(m.get('sim', 0) for m in structured['localMatches']),
        "External": sum(m.get('sim', 0) for m in structured['externalMatches']),
        "AI": sum(m.get('sim', 0) for m in structured['aiMatches'])
    }
    pie_colors = {"Local (Students)": "#ffeb3b", "External": "#ffcdd2", "AI": "#bbdefb"}
    structured['pieData'] = pie_data
    structured['pieColors'] = pie_colors
    return structured

def fetch_external_text(url, timeout=8):
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; PlagiarismChecker/1.0)"}
        r = requests.get(url, headers=headers, timeout=timeout)
        r.raise_for_status()
        html = r.text
        # strip scripts/styles and HTML tags
        text = re.sub(r'(?is)<(script|style).*?>.*?</\1>', '', html)
        text = re.sub(r'(?is)<.*?>', ' ', text)
        text = re.sub(r'\s+', ' ', unescape(text)).strip()
        # truncate to reasonable size
        return text[:200000]
    except Exception:
        logger.exception(f"Failed to fetch external url: {url}")
        return ""

def call_gemini_for_plagiarism(prompt, gemini_api_key, max_attempts=3):
    if not gemini_api_key:
        return None, "No API key provided."
    gemini_api_url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent?key={gemini_api_key}"
    headers = {'Content-Type': 'application/json'}
    payload = {
        "contents": [{"role": "user", "parts": [{"text": prompt}]}],
        "generationConfig": {"responseMimeType": "text/plain"}
    }
    backoff = 1.0
    for attempt in range(1, max_attempts+1):
        try:
            logger.info(f"Calling Gemini (attempt {attempt})...")
            resp = requests.post(gemini_api_url, json=payload, headers=headers, timeout=30)
            resp.raise_for_status()
            jr = resp.json()
            if jr.get('candidates') and jr['candidates'][0].get('content') and jr['candidates'][0]['content'].get('parts') and jr['candidates'][0]['content']['parts'][0].get('text'):
                text = jr['candidates'][0]['content']['parts'][0].get('text')
                return text, None
            else:
                logger.warning("Gemini returned unexpected structure.")
                return None, "Unexpected Gemini response structure."
        except requests.exceptions.HTTPError as he:
            status = getattr(he.response, 'status_code', None)
            logger.warning(f"HTTPError from Gemini: {he} (status {status})")
            if status == 429:
                time.sleep(backoff); backoff *= 2; continue
            else:
                time.sleep(backoff); backoff *= 2; continue
        except requests.exceptions.RequestException as rexc:
            logger.warning(f"RequestException calling Gemini: {rexc}")
            time.sleep(backoff); backoff *= 2; continue
        except Exception:
            logger.exception("Unexpected error calling Gemini")
            time.sleep(backoff); backoff *= 2; continue
    return None, "Retries exhausted or failed to call Gemini."

def build_gemini_prompt(submission_rollno, submitted_document, other_student_sources, external_sources_texts):
    def short_block_for_student(s):
        c = (s.get('content','') or '')[:4000]
        return f"Source Type: Student Paper | ID: {s.get('id')} | Content: {c}"

    student_block = "\n\n".join([short_block_for_student(s) for s in other_student_sources]) or "No reference documents provided."

    external_block_lines = []
    for name, text in external_sources_texts.items():
        external_block_lines.append(f"Source Type: External Source | ID: {name} | Content: {(text or '')[:5000]}")
    external_block = "\n\n".join(external_block_lines) or "No external sources provided."

    prompt = f"""
You are an automated plagiarism detection assistant. Produce a machine-readable JSON object (and nothing else) that compares a Submitted Document against local student papers and external sources.

Header:
Submitted Student ID: {submission_rollno}

Requirements - OUTPUT MUST BE JSON (valid, parseable):
{{
  "submittedRollNo": "<rollno>",
  "overallSimilarity": "<integer percent - overall>",
  "localMatches": [ {{ "sourceId": "Roll No. 1001", "sim": 40 }}, ... ],
  "externalMatches": [ {{ "sourceId": "GeeksforGeeks", "sim": 55 }}, ... ],
  "pieData": {{ "Local (Students)": 40, "External": 50, "AI": 10 }},
  "pieColors": {{ "Local (Students)": "#ffeb3b", "External": "#ffcdd2", "AI": "#bbdefb" }}
}}

Instructions:
- Compare the submitted document (literal text below) to each student paper and each external source.
- Return percentage integers for similarity (0-100).
- For localMatches include only entries where similarity >= 10%.
- For externalMatches include only entries where similarity >= 10%.
- Keep JSON minimal and valid (no extra commentary outside JSON).
- If no matches, return empty arrays for localMatches/externalMatches and pieData values as 0.

Submitted Document:
\"\"\"{(submitted_document or '')[:20000]}\"\"\"


Reference Student Papers:
\"\"\"{student_block}\"\"\"


External Sources:
\"\"\"{external_block}\"\"\"


"""
    return prompt

# default external sources (names only; urls used internally, never returned)
DEFAULT_EXTERNAL_SOURCES = {
    "GeeksforGeeks": "https://www.geeksforgeeks.org/",
    "JavatPoint": "https://www.javatpoint.com/",
    "Programiz": "https://www.programiz.com/",
    "Tutorialspoint": "https://www.tutorialspoint.com/",
    "StackOverflow": "https://stackoverflow.com/",
    "Codecademy": "https://www.codecademy.com/",
    "Unacademy": "https://unacademy.com/",
    "W3Schools": "https://www.w3schools.com/"
}

# ---- Core reusable analysis function ----
def analyze_submission(teacher_id, submitted_content, assignment_id, submission_id, submission_rollno='Unknown', external_urls=None):
    """
    Returns a dict with structured analysis:
    {
      "submittedRollNo": "...",
      "overallSimilarity": "42%",
      "localSummary": {...},   # roll -> %
      "externalSummary": {...},# sourceName -> %
      "pieData": {...},
      "pieColors": {...},
      "pieChartImage": "data:...base64,...",
      "highlightedContent": "<p>...<mark>...</mark>...</p>",
      "geminiUsed": True/False,
      "geminiError": "..." or None
    }
    """
    try:
        # validate teacher
        teacher = get_user_by_id(teacher_id)
        if not teacher or teacher.get('role') != 'teacher':
            return {"error": "Unauthorized: Only teachers can request analysis."}, 403

        # gather other student submissions for the assignment
        other_subs = []
        cursor = submissions_collection.find({"assignmentId": assignment_id, "_id": {"$ne": ObjectId(submission_id)}})
        for s in cursor:
            other_subs.append({"id": f"Roll No. {s.get('rollNo','N/A')}", "roll": s.get('rollNo','N/A'), "content": s.get('content','')})

        # prepare external sources
        external_texts = {}
        # if explicit external_urls passed, use them as External-1, External-2...
        if external_urls and isinstance(external_urls, list) and len(external_urls) > 0:
            for idx, url in enumerate(external_urls):
                name = f"External-{idx+1}"
                u = (url or "").strip()
                if u:
                    external_texts[name] = fetch_external_text(u)
        else:
            # fetch default external sources
            for name, url in DEFAULT_EXTERNAL_SOURCES.items():
                external_texts[name] = fetch_external_text(url) if url else ""

        # Attempt Gemini first (if key available)
        gemini_api_key = os.environ.get("GEMINI_API_KEY", "")
        parsed_structured = None
        gemini_used = False
        gemini_error = None

        if gemini_api_key:
            prompt = build_gemini_prompt(submission_rollno, submitted_content, other_subs, external_texts)
            gemini_raw, gem_err = call_gemini_for_plagiarism(prompt, gemini_api_key, max_attempts=3)
            gemini_error = gem_err
            if gemini_raw:
                # attempt to extract first {...} JSON block and parse
                m = re.search(r'(\{[\s\S]*\})', gemini_raw)
                if m:
                    try:
                        parsed_structured = json.loads(m.group(1))
                        gemini_used = True
                        logger.info("Parsed JSON from Gemini.")
                    except Exception:
                        logger.exception("Failed to parse JSON returned by Gemini.")
                        parsed_structured = None
                else:
                    logger.warning("Gemini response did not include JSON block.")
                    parsed_structured = None
            else:
                logger.warning("Gemini call failed or returned empty; falling back.")

        # Fallback to local heuristic if needed
        if parsed_structured:
            local_matches = parsed_structured.get('localMatches', [])
            external_matches = parsed_structured.get('externalMatches', [])
            pie_data = parsed_structured.get('pieData', {"Local (Students)":0, "External":0, "AI":0})
            pie_colors = parsed_structured.get('pieColors', {"Local (Students)":"#ffeb3b","External":"#ffcdd2","AI":"#bbdefb"})
            overall_similarity = str(parsed_structured.get('overallSimilarity', 0)) + "%"
            # parsed_structured probably does not include phrases; we'll still generate highlightedContent by local sentence-sim checks.
            highlight_source = "gemini"
        else:
            analysis = analyze_with_local_heuristic(submitted_content, [{"roll": s['roll'], "content": s['content']} for s in other_subs], external_texts)
            local_matches = analysis.get('localMatches', [])
            external_matches = analysis.get('externalMatches', [])
            pie_data = analysis.get('pieData', {"Local (Students)":0, "External":0, "AI":0})
            pie_colors = analysis.get('pieColors', {"Local (Students)":"#ffeb3b","External":"#ffcdd2","AI":"#bbdefb"})
            overall_similarity = f"{analysis.get('overallSimilarity',0)}%"
            highlight_source = "local"

        # summaries (ids -> %)
        local_summary, external_summary = summarize_matches_to_ids(local_matches, external_matches)

        pie_img = generate_pie_image(pie_data, pie_colors)

        # Build highlightedContent: mark sentences that match either local or external by re-checking sentences.
        def build_highlighted_html(submitted_text, local_matches_list, external_matches_list, other_subs_list, external_texts_map):
            sentences = safe_sentence_split(submitted_text)
            out_parts = []
            for s in sentences:
                marked = False
                # check local match list first (phrases if present)
                for lm in local_matches_list:
                    # if phrase provided by heuristic, match exact phrase
                    phrase = lm.get('phrase')
                    if phrase and phrase.strip() and phrase.strip() == s.strip():
                        marked = True; break
                if not marked:
                    # fallback: compute similarity with each other_sub
                    for osub in other_subs_list:
                        sim = compute_similarity_percent(s, osub.get('content',''))
                        if sim >= 15:
                            marked = True; break
                if not marked:
                    # check external matches / texts
                    for name, ext_txt in external_texts_map.items():
                        sim = compute_similarity_percent(s, ext_txt)
                        if sim >= 12:
                            marked = True; break
                if marked:
                    out_parts.append(f"<mark style='background:#ffeb3b;padding:0.1rem 0.2rem;border-radius:3px'>{escape_html(s)}</mark>")
                else:
                    out_parts.append(escape_html(s))
            return "<p style='line-height:1.6;font-family:system-ui,Segoe UI,Roboto,Helvetica,Arial, sans-serif;'>" + " ".join(out_parts) + "</p>"

        # helper to escape text for HTML
        def escape_html(t):
            return (t or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

        highlighted_html = build_highlighted_html(submitted_content, local_matches, external_matches, other_subs, external_texts)

        structured = {
            "submittedRollNo": submission_rollno,
            "overallSimilarity": overall_similarity,
            "localSummary": local_summary,
            "externalSummary": external_summary,
            "pieData": pie_data,
            "pieColors": pie_colors,
            "pieChartImage": pie_img,
            "highlightedContent": highlighted_html,
            "geminiUsed": gemini_used,
            "geminiError": gemini_error
        }
        return structured, None
    except Exception as e:
        logger.exception("Error in analyze_submission")
        return None, str(e)

# ----------------- Endpoints -----------------

@app.route('/api/auth/signup', methods=['POST'])
def signup():
    data = request.json
    email = data.get('email')
    password = data.get('password')
    role = data.get('role', 'student')
    branch = data.get('branch')
    photo = data.get('photo')
    if not email or not password or not role:
        return jsonify({"error":"Email, password, and role are required"}), 400
    if not branch:
        return jsonify({"error":"Branch is required for registration"}), 400
    if role == 'student' and not photo:
        return jsonify({"error":"Photo is required for student signup"}), 400
    if users_collection.find_one({"email": email}):
        return jsonify({"error":"User with this email already exists"}), 409
    try:
        user_data = {"email": email, "password": password, "role": role, "branch": branch, "photo": photo if role=='student' else None, "createdAt": datetime.utcnow()}
        result = users_collection.insert_one(user_data)
        return jsonify({"message":"User created successfully","userId":str(result.inserted_id),"userRole":role}), 201
    except Exception:
        logger.exception("Signup error")
        return jsonify({"error":"Internal server error during signup"}), 500

@app.route('/api/auth/login', methods=['POST'])
def login():
    data = request.json
    email = data.get('email')
    password = data.get('password')
    if not email or not password:
        return jsonify({"error":"Email and password are required"}), 400
    user = users_collection.find_one({"email": email})
    if not user or user.get('password') != password:
        return jsonify({"error":"Invalid email or password"}), 401
    return jsonify({"message":"Logged in successfully","userId":str(user['_id']),"userRole":user.get('role')}), 200

@app.route('/api/user/<string:user_id>', methods=['GET'])
def get_user(user_id):
    user = users_collection.find_one({"_id": ObjectId(user_id)})
    if not user:
        return jsonify({"error":"User not found"}), 404
    return jsonify(serialize_doc(user)), 200

@app.route('/api/assignments', methods=['POST'])
def create_assignment():
    data = request.json
    user_id = data.get('userId')
    name = data.get('name')
    description = data.get('description')
    deadline_str = data.get('deadline')
    branches = data.get('branches')
    total_students = data.get('totalStudents')
    user = get_user_by_id(user_id)
    if not user or user['role'] != 'teacher':
        return jsonify({"error":"Unauthorized: Only teachers can create assignments"}), 403
    if not name or not deadline_str or not branches or not isinstance(branches, list) or len(branches)==0 or not total_students:
        return jsonify({"error":"Assignment name, description, deadline, branches and total students are required"}), 400
    try:
        deadline = datetime.fromisoformat(deadline_str.replace('Z','+00:00'))
        assignment_data = {"teacherId": user_id, "name": name, "description": description, "deadline": deadline, "branches": branches, "totalStudents": int(total_students), "createdAt": datetime.utcnow()}
        result = assignments_collection.insert_one(assignment_data)
        return jsonify({"message":"Assignment created successfully","assignmentId":str(result.inserted_id)}), 201
    except Exception:
        logger.exception("Create assignment error")
        return jsonify({"error":"Internal server error creating assignment"}), 500

@app.route('/api/assignments/<string:user_id>', methods=['GET'])
def get_assignments(user_id):
    user = get_user_by_id(user_id)
    if not user:
        return jsonify({"error":"User not found"}), 404
    assignments = []
    if user['role'] == 'teacher':
        cursor = assignments_collection.find({"teacherId": user_id}).sort("createdAt",-1)
    else:
        cursor = assignments_collection.find({"branches": user.get('branch'), "deadline": {"$gt": datetime.utcnow()}}).sort("deadline",1)
    for doc in cursor:
        assignments.append(serialize_doc(doc))
    return jsonify(assignments), 200

@app.route('/api/submissions', methods=['POST'])
def submit_assignment():
    try:
        content = ""
        student_id = (request.form.get('studentId') or (request.json and request.json.get('studentId')))
        assignment_id = (request.form.get('assignmentId') or (request.json and request.json.get('assignmentId')))
        student_name = (request.form.get('name') or (request.json and request.json.get('name')))
        student_rollno = (request.form.get('rollNo') or (request.json and request.json.get('rollNo')))
        if 'file' in request.files and request.files['file'].filename:
            file = request.files['file']
            filename = file.filename.lower()
            file_stream = BytesIO(file.read()); file_stream.seek(0)
            if filename.endswith('.txt'):
                content = file_stream.getvalue().decode('utf-8')
            elif filename.endswith('.pdf'):
                content = extract_text_from_pdf(file_stream)
            elif filename.endswith('.docx'):
                content = extract_text_from_docx(file_stream)
            elif filename.endswith(('.ppt', '.pptx')):
                content = extract_text_from_pptx(file_stream)
            else:
                return jsonify({"error":"Unsupported file type."}), 400
        elif request.json and 'content' in request.json:
            content = request.json.get('content')
        if not student_id or not assignment_id or not content or not student_name or not student_rollno:
            return jsonify({"error":"Missing required data (ID, Name, Roll No, Assignment ID, or Content) for submission"}), 400
        user = users_collection.find_one({"_id": ObjectId(student_id)})
        if not user or user.get('role') != 'student':
            return jsonify({"error":"Unauthorized: Only students can submit assignments"}), 403
        assignment = assignments_collection.find_one({"_id": ObjectId(assignment_id)})
        if not assignment:
            return jsonify({"error":"Assignment not found"}), 404
        if assignment['deadline'] < datetime.utcnow():
            return jsonify({"error":"Assignment deadline has passed"}), 400
        existing = submissions_collection.find_one({"studentId": student_id, "assignmentId": assignment_id})
        if existing:
            return jsonify({"error":"You have already submitted for this assignment"}), 409
        submission_data = {"assignmentId": assignment_id, "studentId": student_id, "name": student_name, "rollNo": student_rollno, "content": content, "submissionDate": datetime.utcnow(), "teacherId": assignment['teacherId']}
        result = submissions_collection.insert_one(submission_data)
        return jsonify({"message":"Assignment submitted successfully","submissionId":str(result.inserted_id)}), 201
    except Exception:
        logger.exception("Submit assignment error")
        return jsonify({"error":"Internal server error during submission"}), 500

@app.route('/api/submissions/<string:assignment_id>', methods=['GET'])
def get_submissions_for_assignment(assignment_id):
    submissions_with_details=[]
    cursor = submissions_collection.find({"assignmentId": assignment_id}).sort("submissionDate", -1)
    cnt=0
    for doc in cursor:
        submissions_with_details.append(serialize_doc(doc)); cnt+=1
    assignment_doc = assignments_collection.find_one({"_id": ObjectId(assignment_id)})
    total_students = assignment_doc.get('totalStudents',0) if assignment_doc else 0
    return jsonify({"submissions": submissions_with_details, "submittedStudentsCount": cnt, "totalStudents": total_students}), 200

@app.route('/api/submissions/student/<string:student_id>', methods=['GET'])
def get_student_submissions(student_id):
    out=[]
    for doc in submissions_collection.find({"studentId": student_id}).sort("submissionDate", -1):
        out.append(serialize_doc(doc))
    return jsonify(out), 200

# ---------------- plagiarism endpoints ----------------

@app.route('/api/plagiarism/check', methods=['POST'])
def check_plagiarism():
    try:
        data = request.json or {}
        teacher_id = data.get('teacherId')
        submitted_content = data.get('submittedContent', '') or ''
        assignment_id = data.get('assignmentId')
        submission_id = data.get('submissionId')
        submission_rollno = data.get('submissionRollNo')
        external_urls = data.get('externalUrls')  # optional list of strings

        # If roll not provided, try to fetch from submission doc
        if not submission_rollno and submission_id:
            try:
                sdoc = submissions_collection.find_one({"_id": ObjectId(submission_id)})
                if sdoc and sdoc.get('rollNo'):
                    submission_rollno = sdoc.get('rollNo')
            except Exception:
                submission_rollno = None
        if not submission_rollno:
            submission_rollno = 'Unknown'

        if not teacher_id or not submitted_content or not assignment_id or not submission_id:
            return jsonify({"error":"Missing required data for plagiarism check"}), 400

        structured, err = analyze_submission(teacher_id, submitted_content, assignment_id, submission_id, submission_rollno=submission_rollno, external_urls=external_urls)
        if structured is None:
            return jsonify({"error": f"Analysis error: {err}"}), 500

        # Build a human readable, but minimal, HTML report for frontend (no copied sentences)
        lines = [f"Submitted Student ID: {structured.get('submittedRollNo','Unknown')}",
                 f"*Overall Similarity Percentage*: {structured.get('overallSimilarity','0%')}",
                 "",
                 "Local Matches (student id : %):"]
        if structured.get('localSummary'):
            for k,v in structured['localSummary'].items():
                lines.append(f"{k}: {v}%")
        else:
            lines.append("None")
        lines.append("")
        lines.append("External Matches (source : %):")
        if structured.get('externalSummary'):
            for k,v in structured['externalSummary'].items():
                lines.append(f"{k}: {v}%")
        else:
            lines.append("None")

        report_text = "\n".join(lines)
        html_report = "<pre style='white-space:pre-wrap;font-family:monospace;background:#fff;padding:10px;border-radius:6px;border:1px solid #e2e8f0;'>" + re.sub(r'&', '&amp;', re.sub(r'<', '&lt;', report_text)) + "</pre>"

        return jsonify({"report": html_report, "structuredReport": structured}), 200
    except Exception:
        logger.exception("Error while processing /api/plagiarism/check")
        return jsonify({"error":"Internal server error during plagiarism check"}), 500

@app.route('/api/plagiarism/check_all', methods=['POST'])
def check_all_submissions():
    """
    Checks plagiarism for ALL submissions for an assignment.
    Expects JSON: { teacherId, assignmentId, externalUrls (optional list) }
    Returns list of student results and flagged students (overallSimilarity >= 50).
    """
    try:
        data = request.json or {}
        teacher_id = data.get('teacherId')
        assignment_id = data.get('assignmentId')
        external_urls = data.get('externalUrls')

        teacher = get_user_by_id(teacher_id)
        if not teacher or teacher.get('role') != 'teacher':
            return jsonify({"error":"Unauthorized: Only teachers can run this."}), 403
        if not assignment_id:
            return jsonify({"error":"assignmentId is required"}), 400

        # fetch all submissions for assignment
        submissions = list(submissions_collection.find({"assignmentId": assignment_id}))
        results = []
        flagged = []

        for s in submissions:
            sid = str(s.get('_id'))
            roll = s.get('rollNo') or 'Unknown'
            student_id = s.get('studentId')
            content = s.get('content', '') or ''

            # analyze each submission
            structured, err = analyze_submission(teacher_id, content, assignment_id, sid, submission_rollno=roll, external_urls=external_urls)
            if structured is None:
                # if error, record it but continue
                results.append({
                    "studentId": student_id,
                    "submissionId": sid,
                    "rollNo": roll,
                    "error": err
                })
                continue

            # overallSimilarity is like "42%"
            try:
                sim_int = int(str(structured.get('overallSimilarity','0%')).replace('%','').strip())
            except Exception:
                sim_int = 0

            entry = {
                "studentId": student_id,
                "submissionId": sid,
                "rollNo": roll,
                "similarity": sim_int,
                "localSummary": structured.get('localSummary', {}),
                "externalSummary": structured.get('externalSummary', {}),
                "pieChartImage": structured.get('pieChartImage'),
                "highlightedContent": structured.get('highlightedContent'),
                "geminiUsed": structured.get('geminiUsed', False)
            }
            results.append(entry)
            if sim_int >= 50:
                flagged.append({"studentId": student_id, "submissionId": sid, "rollNo": roll, "similarity": sim_int})

        return jsonify({"results": results, "flagged": flagged}), 200
    except Exception:
        logger.exception("Error in /api/plagiarism/check_all")
        return jsonify({"error":"Internal server error during check_all"}), 500

@app.route('/api/plagiarism/download', methods=['POST'])
def download_plagiarism_report():
    data = request.json
    report_content = data.get('reportContent')
    filename = data.get('filename', 'plagiarism_report.txt')
    if not report_content:
        return jsonify({"error":"Report content is missing"}), 400
    response = make_response(report_content)
    response.headers['Content-Disposition'] = f'attachment; filename={filename}'
    response.mimetype = 'text/plain'
    return response

@app.route('/')
def index():
    return jsonify({"message":"Plagiarism Detection Backend running"}), 200

if __name__ == '__main__':
    # port 5001 (same as before)
    app.run(host='0.0.0.0', port=5001, debug=True)

